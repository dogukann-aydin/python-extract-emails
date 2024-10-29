import re
import sys
import os
from PyQt6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QPushButton, QFileDialog, QLabel
)
from PyPDF2 import PdfReader
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation

class EmailExtractorApp(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("E-Mail Adresi Çıkarıcı")
        self.setGeometry(100, 100, 400, 300)

        self.layout = QVBoxLayout()

        self.label = QLabel("E-Mail adreslerini çıkarmak için bir veya daha fazla dosya seçin.")
        self.layout.addWidget(self.label)

        self.button_load = QPushButton("Dosyaları Yükle")
        self.button_load.clicked.connect(self.load_files)
        self.layout.addWidget(self.button_load)

        self.setLayout(self.layout)

    def load_files(self):
        file_names, _ = QFileDialog.getOpenFileNames(self, "Dosya Seç", "", "Supported Files (*.txt *.pdf *.xls *.xlsx *.doc *.docx *.ppt *.pptx);;All Files (*)")
        content = ""
        for file_name in file_names:
            if file_name.endswith('.pdf'):
                content += self.read_pdf(file_name) + "\n"
            elif file_name.endswith(('.xls', '.xlsx')):
                content += self.read_excel(file_name) + "\n"
            elif file_name.endswith(('.doc', '.docx')):
                content += self.read_word(file_name) + "\n"
            elif file_name.endswith(('.ppt', '.pptx')):
                content += self.read_ppt(file_name) + "\n"
            elif file_name.endswith('.txt'):
                with open(file_name, 'r', encoding='utf-8') as file:
                    content += file.read() + "\n"

        self.extract_emails(content)

    def read_pdf(self, file_name):
        content = ""
        with open(file_name, 'rb') as file:
            reader = PdfReader(file)
            for page in reader.pages:
                content += page.extract_text() + "\n"
        return content

    def read_excel(self, file_name):
        content = ""
        workbook = load_workbook(filename=file_name, data_only=True)
        for sheet in workbook.sheetnames:
            worksheet = workbook[sheet]
            for row in worksheet.iter_rows(values_only=True):
                content += " ".join(str(cell) for cell in row if cell is not None) + "\n"
        return content

    def read_word(self, file_name):
        content = ""
        doc = Document(file_name)
        for para in doc.paragraphs:
            content += para.text + "\n"
        return content

    def read_ppt(self, file_name):
        content = ""
        presentation = Presentation(file_name)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    content += shape.text + "\n"
        return content

    def extract_emails(self, text):
        emails = re.findall(r'[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}', text)
        
        if emails:
            output_directory = QFileDialog.getExistingDirectory(self, "Çıktı Dizini Seç")
            if output_directory:
                output_file_name = os.path.join(output_directory, "extracted_emails.txt")
                with open(output_file_name, 'w', encoding='utf-8') as output_file:
                    for email in emails:
                        output_file.write(email + '\n')
                self.label.setText(f"{len(emails)} e-posta adresi çıkarıldı ve '{output_file_name}' dosyasına kaydedildi.")
            else:
                self.label.setText("Dizin seçilmedi.")
        else:
            self.label.setText("E-posta adresi bulunamadı.")

if __name__ == "__main__":
    app = QApplication(sys.argv)
    window = EmailExtractorApp()
    window.show()
    sys.exit(app.exec())
